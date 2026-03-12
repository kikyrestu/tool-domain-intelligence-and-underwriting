"""
Crawl Service — fetch source page, extract outbound links, check domain
liveness (DNS + HTTP root), filter domains, and save candidates to DB.
"""

import asyncio
import io
import logging
import os
import random
import re
import time
from datetime import datetime, timezone
from urllib.parse import urljoin, urlparse, parse_qs, unquote

import dns.asyncresolver
import dns.exception
import httpx
from bs4 import BeautifulSoup
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.source import Source
from app.models.crawl_job import CrawlJob
from app.models.candidate import CandidateDomain
from app.services.proxy_service import proxy_service
from app.services.toxicity_service import TOXICITY_PATTERNS
from app.utils.ssrf_guard import is_safe_url
from app.utils.domain_filter import extract_domain, is_valid_candidate, BLACKLIST
from app.config import get_settings

logger = logging.getLogger(__name__)

# Compiled parking patterns reused across requests
_PARKING_RE = re.compile(
    "|".join(TOXICITY_PATTERNS["parking"]), re.IGNORECASE
)


def _is_parked(html: str) -> bool:
    """Return True if the HTTP response body looks like a domain parking page."""
    if not html:
        return False
    # Only scan the first 8 KB — parking indicators are always in <head>/<body> preamble
    return bool(_PARKING_RE.search(html[:8192]))


TIMEOUT = 20
MAX_CONCURRENT = 5

_CHALLENGE_RE = re.compile(
    r"(enable javascript and cookies|just a moment\.\.\.?|one moment, please|cf-browser-verification|"
    r"challenge-platform|please wait\.\.\.|ddos protection by cloudflare|"
    r"checking your browser|are you a robot|access denied|ray id|"
    r"unsupported media type|403 forbidden|503 service)",
    re.IGNORECASE,
)

# HTTP error status codes that should trigger fallback regardless of body content
_BAD_STATUSES = {400, 403, 415, 429, 503}


def _is_challenge_page(html: str, status: int = 200) -> bool:
    """Return True if response looks like a challenge/error page, not real content."""
    if status in _BAD_STATUSES:
        return True
    return len(html) < 60_000 and bool(_CHALLENGE_RE.search(html))

USER_AGENTS = [
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 14_3) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.2 Safari/605.1.15",
    "Mozilla/5.0 (X11; Linux x86_64; rv:123.0) Gecko/20100101 Firefox/123.0",
]


def _get_headers() -> dict:
    return {
        "User-Agent": random.choice(USER_AGENTS),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
        "Upgrade-Insecure-Requests": "1",
        "Sec-Fetch-Dest": "document",
        "Sec-Fetch-Mode": "navigate",
        "Sec-Fetch-Site": "none",
        "Sec-Fetch-User": "?1",
        "Cache-Control": "max-age=0",
    }


def _make_client() -> httpx.AsyncClient:
    """Create an httpx client with optional proxy."""
    proxy = proxy_service.get_next()
    return httpx.AsyncClient(proxy=proxy, headers=_get_headers(), follow_redirects=True)


# ---------------------------------------------------------------------------
# Multi-provider scraping API pool — round-robin per provider + per key
# ---------------------------------------------------------------------------
_provider_key_indices: dict[str, int] = {}


def _parse_keys(raw: str) -> list[str]:
    """Split comma-separated key string into list."""
    return [k.strip() for k in raw.split(",") if k.strip()]


def _next_key(provider: str, keys: list[str]) -> str:
    """Round-robin key selection per provider."""
    idx = _provider_key_indices.get(provider, 0)
    key = keys[idx % len(keys)]
    _provider_key_indices[provider] = idx + 1
    return key


def _zenrows_key_pool(settings) -> list[str]:
    """Return ZenRows keys (ZENROWS_API_KEYS preferred, ZENROWS_API_KEY fallback)."""
    if settings.ZENROWS_API_KEYS:
        keys = _parse_keys(settings.ZENROWS_API_KEYS)
        if keys:
            return keys
    if settings.ZENROWS_API_KEY:
        return [settings.ZENROWS_API_KEY]
    return []


async def _try_scraperapi(url: str, key: str, js: bool = False, residential: bool = False) -> str | None:
    """Call ScraperAPI. js=True costs more credits. residential costs extra."""
    params = {"api_key": key, "url": url, "keep_headers": "true"}
    if js:
        params["render"] = "true"
    if residential:
        params["premium"] = "true"  # residential proxy tier
    try:
        async with httpx.AsyncClient(follow_redirects=True) as client:
            resp = await client.get("https://api.scraperapi.com/", params=params, timeout=60)
            if resp.status_code == 200 and len(resp.text) > 500 and not _is_challenge_page(resp.text, resp.status_code):
                return resp.text
    except Exception as e:
        logger.debug("ScraperAPI error: %s", e)
    return None


async def _try_scrapingbee(url: str, key: str, js: bool = False, stealth: bool = False) -> str | None:
    """Call Scrapingbee. js=True costs 5 credits. stealth costs 10 credits."""
    params = {"api_key": key, "url": url, "block_resources": "false"}
    if stealth:
        params["stealth_proxy"] = "true"
    elif js:
        params["render_js"] = "true"
    try:
        async with httpx.AsyncClient(follow_redirects=True) as client:
            resp = await client.get("https://app.scrapingbee.com/api/v1/", params=params, timeout=60)
            if resp.status_code == 200 and len(resp.text) > 500 and not _is_challenge_page(resp.text, resp.status_code):
                return resp.text
    except Exception as e:
        logger.debug("Scrapingbee error: %s", e)
    return None


async def _try_crawlbase(url: str, key: str, js: bool = False) -> str | None:
    """Call Crawlbase. js=False uses normal token (1cr), js=True uses JS token (5cr, headless)."""
    params = {"token": key, "url": url}
    if js:
        params["ajax_wait"] = "true"
        params["page_wait"] = "3000"
    try:
        async with httpx.AsyncClient(follow_redirects=True) as client:
            resp = await client.get("https://api.crawlbase.com/", params=params, timeout=60)
            if resp.status_code == 200 and len(resp.text) > 500 and not _is_challenge_page(resp.text, resp.status_code):
                return resp.text
    except Exception as e:
        logger.debug("Crawlbase error: %s", e)
    return None


async def _try_scrapegraphai(url: str, key: str) -> str | None:
    """Call ScrapeGraphAI to fetch raw HTML of a page (markdownify endpoint)."""
    try:
        async with httpx.AsyncClient(follow_redirects=True) as client:
            resp = await client.post(
                "https://api.scrapegraphai.com/v1/markdownify",
                json={"website_url": url},
                headers={"SGAI-APIKEY": key, "Content-Type": "application/json"},
                timeout=60,
            )
            if resp.status_code == 200:
                data = resp.json()
                # Returns {"result": "<markdown text>"}
                content = data.get("result", "")
                if content and len(content) > 300:
                    return content
    except Exception as e:
        logger.debug("ScrapeGraphAI error: %s", e)
    return None


async def _fetch_page(url: str) -> str | None:
    """Fetch a page.
    Strategy:
    - If ZenRows key(s) set → rotate keys round-robin, tiered: basic → js_render → antibot
    - Otherwise: direct → proxy → cloudscraper → Playwright+Stealth → 2captcha
    """
    settings = get_settings()

    # ── Provider pools ────────────────────────────────────────────────────
    zr_keys    = _zenrows_key_pool(settings)
    sa_keys    = _parse_keys(settings.SCRAPERAPI_KEYS)  if settings.SCRAPERAPI_KEYS  else []
    bee_keys   = _parse_keys(settings.SCRAPINGBEE_KEYS) if settings.SCRAPINGBEE_KEYS else []
    cb_keys    = _parse_keys(settings.CRAWLBASE_KEYS)   if settings.CRAWLBASE_KEYS   else []
    cb_js_keys = _parse_keys(settings.CRAWLBASE_JS_KEYS) if settings.CRAWLBASE_JS_KEYS else []
    sgai_key   = settings.SCRAPEGRAPHAI_KEY or None

    # ── Provider 1: ZenRows — basic → js_render → antibot ────────────────
    if zr_keys:
        try:
            from zenrows import ZenRowsClient
            for zr_params, label in [
                ({}, "basic"),
                ({"js_render": "true"}, "js_render"),
                ({"js_render": "true", "antibot": "true"}, "antibot"),
            ]:
                api_key = _next_key("zenrows", zr_keys)
                client_zr = ZenRowsClient(api_key)
                resp_zr = await client_zr.get_async(url, params=zr_params)
                if resp_zr and len(resp_zr.text) > 500 and not _is_challenge_page(resp_zr.text, resp_zr.status_code):
                    logger.info("ZenRows[%s] ...%s OK %s (%d chars)", label, api_key[-6:], url, len(resp_zr.text))
                    return resp_zr.text
                logger.debug("ZenRows[%s] bad response %s, escalating", label, url)
        except Exception as e:
            logger.debug("ZenRows failed %s: %s", url, e)

    # ── Provider 2: ScraperAPI — plain → js → residential ────────────────
    if sa_keys:
        for js, residential, label in [(False, False, "plain"), (True, False, "js"), (True, True, "residential")]:
            key = _next_key("scraperapi", sa_keys)
            result = await _try_scraperapi(url, key, js=js, residential=residential)
            if result:
                logger.info("ScraperAPI[%s] ...%s OK %s (%d chars)", label, key[-6:], url, len(result))
                return result
            logger.debug("ScraperAPI[%s] bad response %s, escalating", label, url)

    # ── Provider 3: Scrapingbee — plain → js → stealth ───────────────────
    if bee_keys:
        for js, stealth, label in [(False, False, "plain"), (True, False, "js"), (False, True, "stealth")]:
            key = _next_key("scrapingbee", bee_keys)
            result = await _try_scrapingbee(url, key, js=js, stealth=stealth)
            if result:
                logger.info("Scrapingbee[%s] ...%s OK %s (%d chars)", label, key[-6:], url, len(result))
                return result
            logger.debug("Scrapingbee[%s] bad response %s, escalating", label, url)

    # ── Provider 4: Crawlbase — normal → JS (headless+residential) ──────────
    if cb_keys:
        key = _next_key("crawlbase", cb_keys)
        result = await _try_crawlbase(url, key, js=False)
        if result:
            logger.info("Crawlbase[normal] ...%s OK %s (%d chars)", key[-6:], url, len(result))
            return result
        logger.debug("Crawlbase[normal] bad response %s, escalating to JS", url)

    if cb_js_keys:
        key = _next_key("crawlbase_js", cb_js_keys)
        result = await _try_crawlbase(url, key, js=True)
        if result:
            logger.info("Crawlbase[js] ...%s OK %s (%d chars)", key[-6:], url, len(result))
            return result
        logger.debug("Crawlbase[js] bad response %s", url)

    # ── Provider 5: ScrapeGraphAI — AI-based last resort ──────────────────
    if sgai_key:
        result = await _try_scrapegraphai(url, sgai_key)
        if result:
            logger.info("ScrapeGraphAI OK %s (%d chars)", url, len(result))
            return result
        logger.debug("ScrapeGraphAI bad response %s", url)

    # Fallback 1: direct HTTP/1.1
    try:
        async with httpx.AsyncClient(headers=_get_headers(),
                                     follow_redirects=True) as client:
            resp = await client.get(url, timeout=TIMEOUT)
            if len(resp.text) > 500 and not _is_challenge_page(resp.text, resp.status_code):
                logger.debug("Direct fetch succeeded for %s (%d chars)", url, len(resp.text))
                return resp.text
            else:
                logger.debug("Direct got bad/challenge response for %s (status=%d, %d chars), trying next",
                             url, resp.status_code, len(resp.text))
    except Exception as e:
        logger.debug("Direct fetch failed for %s: %s, trying proxy", url, e)

    # Fallback 2: proxy
    proxy = proxy_service.get_next()
    if proxy:
        try:
            async with httpx.AsyncClient(proxy=proxy, headers=_get_headers(),
                                         follow_redirects=True) as client:
                resp = await client.get(url, timeout=TIMEOUT)
                if len(resp.text) > 500 and not _is_challenge_page(resp.text, resp.status_code):
                    logger.debug("Proxy fetch succeeded for %s (%d chars)", url, len(resp.text))
                    return resp.text
                else:
                    logger.debug("Proxy got bad/challenge response for %s (status=%d), trying cloudscraper",
                                 url, resp.status_code)
        except Exception as e:
            logger.debug("Proxy also failed for %s: %s", url, e)

    # Fallback 3: cloudscraper
    try:
        import cloudscraper
        import functools
        scraper = cloudscraper.create_scraper()
        loop = asyncio.get_event_loop()
        resp_cs = await loop.run_in_executor(
            None,
            functools.partial(scraper.get, url, timeout=TIMEOUT),
        )
        if resp_cs and len(resp_cs.text) > 500 and not _is_challenge_page(resp_cs.text, resp_cs.status_code):
            logger.debug("cloudscraper succeeded for %s (%d chars)", url, len(resp_cs.text))
            return resp_cs.text
        elif resp_cs:
            logger.debug("cloudscraper got challenge page for %s (status=%d, %d chars), trying Playwright",
                         url, resp_cs.status_code, len(resp_cs.text))
    except Exception as e:
        logger.debug("cloudscraper failed for %s: %s, trying Playwright", url, e)

    # Fallback 4: Playwright+Stealth headless browser
    try:
        from playwright.async_api import async_playwright
        from playwright_stealth import Stealth
        async with Stealth().use_async(async_playwright()) as pw:
            browser = await pw.chromium.launch(headless=True)
            context = await browser.new_context(
                user_agent=random.choice(USER_AGENTS),
                locale="en-US",
            )
            page = await context.new_page()
            await page.goto(url, wait_until="domcontentloaded", timeout=TIMEOUT * 1000)
            await page.wait_for_timeout(5000)
            html = await page.content()
            await browser.close()
            if html and len(html) > 500 and not _is_challenge_page(html):
                logger.info("Playwright+Stealth succeeded for %s (%d chars)", url, len(html))
                return html
            else:
                logger.debug("Playwright+Stealth still got challenge page for %s (%d chars)", url, len(html) if html else 0)
    except Exception as e:
        logger.debug("Playwright+Stealth failed for %s: %s", url, e)

    # Fallback 5: Playwright + 2captcha Turnstile solver
    if settings.TWOCAPTCHA_API_KEY:
        try:
            from playwright.async_api import async_playwright
            from playwright_stealth import Stealth
            from twocaptcha import TwoCaptcha
            import functools

            async with Stealth().use_async(async_playwright()) as pw:
                browser = await pw.chromium.launch(headless=True)
                context = await browser.new_context(
                    user_agent=random.choice(USER_AGENTS),
                    locale="en-US",
                )
                page = await context.new_page()
                await page.goto(url, wait_until="domcontentloaded", timeout=TIMEOUT * 1000)
                await page.wait_for_timeout(3000)

                # Extract Turnstile sitekey
                sitekey = await page.evaluate("""
                    () => {
                        const el = document.querySelector('[data-sitekey]');
                        return el ? el.getAttribute('data-sitekey') : null;
                    }
                """)

                if sitekey:
                    logger.info("Solving Cloudflare Turnstile via 2captcha for %s (sitekey=%s)", url, sitekey)
                    solver = TwoCaptcha(settings.TWOCAPTCHA_API_KEY)
                    loop = asyncio.get_event_loop()
                    result = await loop.run_in_executor(
                        None,
                        functools.partial(solver.turnstile, sitekey=sitekey, url=url)
                    )
                    token = result.get("code") if result else None

                    if token:
                        # Inject token and submit
                        await page.evaluate("""
                            (token) => {
                                const field = document.querySelector('[name="cf-turnstile-response"]');
                                if (field) field.value = token;
                                // Also set via window callback if available
                                const widget = document.querySelector('.cf-turnstile');
                                if (widget && widget.__CF_CHLOUT__) widget.__CF_CHLOUT__(token);
                            }
                        """, token)
                        await page.wait_for_timeout(2000)
                        # Re-grab content after solving
                        html = await page.content()
                        await browser.close()
                        if html and len(html) > 500 and not _is_challenge_page(html):
                            logger.info("2captcha Turnstile solved for %s (%d chars)", url, len(html))
                            return html
                else:
                    await browser.close()
        except Exception as e:
            logger.debug("2captcha fallback failed for %s: %s", url, e)
    else:
        logger.debug("No TWOCAPTCHA_API_KEY set, skipping 2captcha fallback")

    logger.error("Failed to fetch %s (all methods failed)", url)
    return None


def _unwrap_redirect(href: str, source_url: str) -> str:
    """
    Unwrap forum redirect/proxy URLs.
    e.g. /proxy.php?link=https://example.com  → https://example.com
         /goto/https://example.com             → https://example.com
    """
    full = urljoin(source_url, href)
    parsed = urlparse(full)

    # Try query params: ?link=, ?url=, ?target=, ?href=, ?to=
    for param in ("link", "url", "target", "href", "to", "u", "q"):
        val = parse_qs(parsed.query).get(param)
        if val:
            candidate = unquote(val[0])
            if candidate.startswith(("http://", "https://")):
                return candidate

    # Try path-based redirect: /goto/https://example.com or /redirect/https://example.com
    path = parsed.path
    match = re.search(r"/(goto|redirect|out|external)[/=](https?://.+)", path)
    if match:
        return unquote(match.group(2))

    return full


# Regex untuk temukan URL plain-text di konten HTML / teks biasa
_URL_RE = re.compile(
    r'(?<!["\'=>])(https?://[a-zA-Z0-9\-._~:/?#\[\]@!$&\'()*+,;=%]+)',
    re.IGNORECASE,
)

# File extensions that indicate a binary/document source URL (not HTML)
_BINARY_EXTENSIONS = {
    "pdf", "pptx", "ppt", "docx", "doc", "xlsx", "xls",
    "odt", "ods", "odp",
}

# File extensions to skip when extracting links from HTML — these are page paths,
# not real domains (e.g. /page.asp, /forum.cfm, /index.php)
_PAGE_EXT_RE = re.compile(
    r"\.(?:asp|aspx|cfm|cfml|jsp|jspx|php|php[3-5]|phtml|rb|py|pl|cgi|do|action)"
    r"(?:[?#]|$)",
    re.IGNORECASE,
)


def _extract_outbound_links(html: str, source_url: str) -> list[tuple[str, str]]:
    """Extract outbound links from HTML. Returns list of (url, domain).
    Tolerant parser: handles legacy HTML (broken tags, relative URLs,
    meta refresh redirects, non-UTF-8 encoding already decoded upstream).
    """
    soup = BeautifulSoup(html, "html.parser")
    source_domain = extract_domain(source_url)

    results = []
    seen_domains = set()

    def _add(url: str):
        parsed_url = urlparse(url)
        # Skip URLs whose path ends with a server-side page extension —
        # these are page paths on an existing domain, not domain candidates.
        if _PAGE_EXT_RE.search(parsed_url.path):
            return
        domain = extract_domain(url)
        if not domain or domain == source_domain:
            return
        if not is_valid_candidate(domain):
            return
        if domain in seen_domains:
            return
        seen_domains.add(domain)
        results.append((url, domain))

    # 1. Anchor tags (dengan unwrap redirect/proxy forum)
    for tag in soup.find_all("a", href=True):
        href = tag["href"].strip()
        if not href or href.startswith(("#", "mailto:", "javascript:")):
            continue
        full_url = _unwrap_redirect(href, source_url)
        parsed = urlparse(full_url)
        if parsed.scheme not in ("http", "https"):
            continue
        _add(full_url)

    # 2. Meta refresh — legacy HTML redirect: <meta http-equiv="refresh" content="0;url=https://...">
    for meta in soup.find_all("meta"):
        http_equiv = (meta.get("http-equiv") or "").lower()
        if http_equiv == "refresh":
            content = meta.get("content", "")
            # Format: "N;url=https://..." or "N;URL=https://..."
            match = re.search(r"url\s*=\s*([^\s;\"']+)", content, re.IGNORECASE)
            if match:
                ref_url = match.group(1).strip().strip("'\"")
                full_url = urljoin(source_url, ref_url)
                if urlparse(full_url).scheme in ("http", "https"):
                    _add(full_url)

    # 3. Plain-text URLs di body (untuk forum post, blog, dll)
    for match in _URL_RE.finditer(html):
        url = match.group(1).rstrip(".,;)'\"")
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            continue
        _add(url)

    return results


# ---------------------------------------------------------------------------
# Binary document fetcher + text extractors
# ---------------------------------------------------------------------------
_MAX_BINARY_SIZE = 25 * 1024 * 1024  # 25 MB hard limit


async def _fetch_binary(url: str) -> bytes | None:
    """Download a binary file (PDF, PPTX, etc.) up to _MAX_BINARY_SIZE."""
    if not is_safe_url(url):
        logger.warning("SSRF guard blocked binary fetch: %s", url)
        return None
    try:
        async with httpx.AsyncClient(headers=_get_headers(), follow_redirects=True) as client:
            async with client.stream("GET", url, timeout=60) as resp:
                if resp.status_code != 200:
                    return None
                content_length = int(resp.headers.get("content-length", 0))
                if content_length and content_length > _MAX_BINARY_SIZE:
                    logger.warning("Binary file too large (%d bytes): %s", content_length, url)
                    return None
                chunks = []
                total = 0
                async for chunk in resp.aiter_bytes(65536):
                    total += len(chunk)
                    if total > _MAX_BINARY_SIZE:
                        logger.warning("Binary file exceeded size limit mid-stream: %s", url)
                        return None
                    chunks.append(chunk)
                return b"".join(chunks)
    except Exception as e:
        logger.debug("Binary fetch failed for %s: %s", url, e)
        return None


def _extract_text_from_pdf(data: bytes) -> str:
    """Extract all text from a PDF binary."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        return "\n".join(text_parts)
    except Exception as e:
        logger.debug("PDF extraction failed: %s", e)
        return ""


def _extract_text_from_pptx(data: bytes) -> str:
    """Extract all text from a PPTX binary."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        logger.debug("PPTX extraction failed: %s", e)
        return ""


def _extract_text_from_docx(data: bytes) -> str:
    """Extract all text from a DOCX binary."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.debug("DOCX extraction failed: %s", e)
        return ""


def _extract_text_from_xlsx(data: bytes) -> str:
    """Extract all text/URLs from an XLSX binary."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        parts.append(str(cell))
        return " ".join(parts)
    except Exception as e:
        logger.debug("XLSX extraction failed: %s", e)
        return ""


def _binary_to_text(data: bytes, ext: str) -> str:
    """Dispatch binary data to the correct text extractor by file extension."""
    ext = ext.lower()
    if ext == "pdf":
        return _extract_text_from_pdf(data)
    elif ext in ("pptx", "ppt"):
        return _extract_text_from_pptx(data)
    elif ext in ("docx", "doc"):
        return _extract_text_from_docx(data)
    elif ext in ("xlsx", "xls"):
        return _extract_text_from_xlsx(data)
    return ""


def _extract_links_from_text(text: str, source_url: str) -> list[tuple[str, str]]:
    """Run URL regex over plain text and return valid (url, domain) tuples."""
    source_domain = extract_domain(source_url)
    results = []
    seen_domains: set[str] = set()
    for match in _URL_RE.finditer(text):
        url = match.group(1).rstrip(".,;)'\"")
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            continue
        if _PAGE_EXT_RE.search(parsed.path):
            continue
        domain = extract_domain(url)
        if not domain or domain == source_domain:
            continue
        if not is_valid_candidate(domain):
            continue
        if domain in seen_domains:
            continue
        seen_domains.add(domain)
        results.append((url, domain))
    return results


async def _check_domain(domain: str, semaphore: asyncio.Semaphore) -> dict:
    """Check if a DOMAIN is alive by checking DNS + HTTP root.
    If domain is dead, immediately check RDAP for buy availability.

    Returns dict with:
      - dns_resolves: bool — does the domain have A/AAAA records?
      - http_status: int|None — status code of GET https://domain/
      - is_domain_alive: bool — combined verdict
      - availability_status: str|None — RDAP result for dead domains
      - whois_registrar: str|None
      - whois_created_date: date|None
      - whois_expiry_date: date|None
      - whois_days_left: int|None
    """
    from app.services.whois_service import _rdap_lookup

    # HTTP status codes that mean a real server responded (domain is registered + active).
    # 403 and 429 are intentionally excluded: they indicate Cloudflare/WAF blocks or
    # rate-limiting, NOT that the domain has live content. Including them caused domains
    # blocked by Cloudflare to be marked "alive", which could then produce a false
    # "Available" label if RDAP also returned an ambiguous result.
    # 401/402/405/407 kept — they indicate an actual web server is running.
    ALIVE_STATUSES = {200, 201, 204, 301, 302, 303, 307, 308, 401, 402, 405, 407}

    # Status codes that indicate a WAF/CDN block — domain has DNS but content is gated.
    # We still mark dns_resolves=True but treat the domain as "alive" only if DNS confirms it.
    BLOCKED_STATUSES = {403, 429}

    result = {"dns_resolves": False, "http_status": None, "is_domain_alive": False,
              "dns_mx_records": False, "is_parked": False,
              "availability_status": None, "whois_registrar": None,
              "whois_created_date": None, "whois_expiry_date": None,
              "whois_days_left": None}

    async with semaphore:
        # --- Layer 1: DNS Resolution ---
        try:
            answers = await dns.asyncresolver.resolve(domain, "A")
            if answers:
                result["dns_resolves"] = True
        except (dns.exception.DNSException, Exception):
            pass

        if not result["dns_resolves"]:
            try:
                answers = await dns.asyncresolver.resolve(domain, "AAAA")
                if answers:
                    result["dns_resolves"] = True
            except (dns.exception.DNSException, Exception):
                pass

        # --- Layer 1b: MX record check (parallel signal of active email use) ---
        try:
            mx_answers = await dns.asyncresolver.resolve(domain, "MX")
            if mx_answers:
                result["dns_mx_records"] = True
        except (dns.exception.DNSException, Exception):
            pass

        # No DNS = domain is truly dead → check if buyable
        if not result["dns_resolves"]:
            rdap = await _rdap_lookup(domain, dns_resolves=False)
            result["availability_status"] = rdap["status"]
            result["whois_registrar"] = rdap.get("registrar")
            result["whois_created_date"] = rdap.get("created_date")
            result["whois_expiry_date"] = rdap.get("expiry_date")
            result["whois_days_left"] = rdap.get("days_left")
            logger.info("  Dead domain %s → RDAP: %s", domain, rdap["status"])
            return result

        # --- Layer 2: HTTP root domain check ---
        root_url = f"https://{domain}/"
        if not is_safe_url(root_url):
            result["is_domain_alive"] = True  # DNS resolves, assume alive
            return result

        # Try HTTPS first
        for url in [f"https://{domain}/", f"http://{domain}/"]:
            try:
                async with httpx.AsyncClient(headers=_get_headers(),
                                             follow_redirects=True, http2=True) as client:
                    resp = await client.get(url, timeout=TIMEOUT)
                    result["http_status"] = resp.status_code
                    if resp.status_code in ALIVE_STATUSES:
                        result["is_domain_alive"] = True
                        # --- Layer 2b: parking detection on live 200 responses ---
                        if resp.status_code == 200:
                            result["is_parked"] = _is_parked(resp.text)
                            if result["is_parked"]:
                                logger.info("  Parked domain detected: %s", domain)
                    elif resp.status_code in BLOCKED_STATUSES:
                        # WAF/CDN block — domain has DNS so it's registered, but
                        # content is inaccessible. Mark alive via DNS, not HTTP.
                        result["is_domain_alive"] = result["dns_resolves"]
                    else:
                        result["is_domain_alive"] = False
                    # --- Layer 3: RDAP for ALL domains (dead AND alive) ---
                    rdap = await _rdap_lookup(domain, dns_resolves=result["dns_resolves"])
                    result["availability_status"] = rdap["status"]
                    result["whois_registrar"] = rdap.get("registrar")
                    result["whois_created_date"] = rdap.get("created_date")
                    result["whois_expiry_date"] = rdap.get("expiry_date")
                    result["whois_days_left"] = rdap.get("days_left")
                    if result["is_domain_alive"]:
                        logger.info("  Alive domain %s (HTTP %s) → RDAP: %s", domain, resp.status_code, rdap["status"])
                    else:
                        logger.info("  Dead domain %s (HTTP %s) → RDAP: %s", domain, resp.status_code, rdap["status"])
                    return result
            except Exception:
                continue

        # DNS resolves but HTTP completely failed → dead, check buyable
        result["is_domain_alive"] = False
        rdap = await _rdap_lookup(domain, dns_resolves=result["dns_resolves"])
        result["availability_status"] = rdap["status"]
        result["whois_registrar"] = rdap.get("registrar")
        result["whois_created_date"] = rdap.get("created_date")
        result["whois_expiry_date"] = rdap.get("expiry_date")
        result["whois_days_left"] = rdap.get("days_left")
        logger.info("  Dead domain %s (HTTP fail) → RDAP: %s", domain, rdap["status"])
        return result


async def run_crawl(source_id: int, db: AsyncSession):
    """
    Main crawl pipeline:
    1. Create CrawlJob
    2. Fetch source page
    3. Extract + filter outbound links
    4. Check dead links
    5. Save candidates to DB
    6. Update CrawlJob stats
    """
    settings = get_settings()

    # Get source
    source = await db.get(Source, source_id)
    if not source:
        raise ValueError(f"Source {source_id} not found")

    # Create crawl job
    job = CrawlJob(source_id=source_id, status="running", current_step="crawling", started_at=datetime.now(timezone.utc))
    db.add(job)
    await db.flush()

    try:
        # 1. Detect source type and route to appropriate pipeline
        _src = source.url
        _path = urlparse(_src).path.lower() if _src.startswith("http") else ""
        source_ext = os.path.splitext(_path)[1].lstrip(".")

        if _src.lower().startswith("crtsh://"):
            # ── Pipeline A: Certificate Transparency (crt.sh) ──────────────
            from app.services.crtsh_service import fetch_domains_from_crtsh, parse_crtsh_keyword
            keyword = parse_crtsh_keyword(_src) or ""
            if not keyword:
                job.status = "failed"
                job.error_message = "crtsh:// URL must include a keyword, e.g. crtsh://technology"
                job.completed_at = datetime.now(timezone.utc)
                await db.commit()
                return job
            job.current_step = "crtsh"
            await db.flush()
            links = await fetch_domains_from_crtsh(keyword)
            logger.info("crt.sh '%s': %d domain candidates", keyword, len(links))

        elif _path.rstrip("/").endswith("robots.txt"):
            # ── Pipeline B: robots.txt → sitemap auto-discovery ────────────
            from app.services.sitemap_service import fetch_links_from_robots
            job.current_step = "robots_txt"
            await db.flush()
            links = await fetch_links_from_robots(_src)
            logger.info("robots.txt '%s': %d domain candidates", _src, len(links))

        elif _path.endswith(".xml") or _path.endswith(".xml.gz") or "sitemap" in _path:
            # ── Pipeline C: Sitemap XML (direct or index) ──────────────────
            from app.services.sitemap_service import fetch_links_from_sitemap
            job.current_step = "sitemap"
            await db.flush()
            links = await fetch_links_from_sitemap(_src)
            logger.info("Sitemap '%s': %d domain candidates", _src, len(links))

        elif source_ext in _BINARY_EXTENSIONS:
            # Binary document pipeline
            logger.info("Binary source detected (%s): %s", source_ext, source.url)
            binary_data = await _fetch_binary(source.url)
            if not binary_data:
                job.status = "failed"
                job.error_message = "Failed to fetch binary source file"
                job.completed_at = datetime.now(timezone.utc)
                await db.commit()
                return job
            text_content = _binary_to_text(binary_data, source_ext)
            if not text_content:
                job.status = "failed"
                job.error_message = f"Could not extract text from {source_ext.upper()} file"
                job.completed_at = datetime.now(timezone.utc)
                await db.commit()
                return job
            links = _extract_links_from_text(text_content, source.url)
            logger.info("Binary (%s): %d bytes → %d chars text → %d links from %s",
                        source_ext, len(binary_data), len(text_content), len(links), source.url)
        else:
            # Standard HTML pipeline
            html = await _fetch_page(source.url)
            if not html:
                job.status = "failed"
                job.error_message = "Failed to fetch source page"
                job.completed_at = datetime.now(timezone.utc)
                await db.commit()
                return job
            links = _extract_outbound_links(html, source.url)
            logger.info("Fetched HTML: %d chars, extracted %d links from %s", len(html), len(links), source.url)

        # Also scan for binary document links embedded in HTML pages and queue text extraction
        # (links to .pdf/.pptx/.docx found in anchor tags are already carried as (url, domain)
        #  by _extract_outbound_links — the domain extraction there is correct, no extra step needed)
        links = links[:settings.MAX_CANDIDATES_PER_CRAWL]

        with db.no_autoflush:
            job.total_links_found = len(links)

        # 3. Check domains concurrently (using ROOT domain, not specific URL)
        semaphore = asyncio.Semaphore(MAX_CONCURRENT)
        # Deduplicate domains for checking (avoid checking same domain twice)
        unique_domains = list({domain for _, domain in links})
        check_tasks = [_check_domain(d, semaphore) for d in unique_domains]
        check_results = await asyncio.gather(*check_tasks)
        domain_status = dict(zip(unique_domains, check_results))

        # 4. Save candidates
        dead_count = 0
        saved = 0
        for (url, domain) in links:
            status = domain_status[domain]
            if not status["is_domain_alive"]:
                dead_count += 1

            # Check if domain already exists for this source
            existing = await db.execute(
                select(CandidateDomain).where(
                    CandidateDomain.domain == domain,
                    CandidateDomain.source_id == source_id,
                )
            )
            if existing.scalar_one_or_none():
                continue

            candidate = CandidateDomain(
                domain=domain,
                source_id=source_id,
                crawl_job_id=job.id,
                source_url_found=source.url,
                original_link=url,
                niche=source.niche,
                http_status=status["http_status"],
                dns_resolves=status["dns_resolves"],
                dns_mx_records=status.get("dns_mx_records", False),
                is_domain_alive=status["is_domain_alive"],
                is_parked=status.get("is_parked", False),
                availability_status=status.get("availability_status"),
                whois_registrar=status.get("whois_registrar"),
                whois_created_date=status.get("whois_created_date"),
                whois_expiry_date=status.get("whois_expiry_date"),
                whois_days_left=status.get("whois_days_left"),
                whois_checked_at=datetime.now(timezone.utc) if status.get("availability_status") else None,
                # Provenance
                source_type=_provenance["source_type"],
                parser_type=_provenance["parser_type"],
                source_origin=_provenance["source_origin"],
                extraction_note=_provenance["extraction_note"],
            )
            db.add(candidate)
            saved += 1

        job.total_candidates = saved
        job.total_dead_links = dead_count
        job.status = "completed"
        job.completed_at = datetime.now(timezone.utc)

        await db.commit()
        logger.info("Crawl completed: %d candidates, %d dead links", saved, dead_count)
        return job

    except Exception as e:
        logger.error("Crawl failed: %s", e)
        try:
            await db.rollback()
            job.status = "failed"
            job.error_message = str(e)[:500]
            job.completed_at = datetime.now(timezone.utc)
            await db.commit()
        except Exception as commit_err:
            logger.error("Failed to save crawl failure state: %s", commit_err)
        return job
